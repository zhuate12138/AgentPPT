"""Preview service - converts PPT to images."""

import asyncio
import subprocess
from pathlib import Path
from typing import List, Optional
import shutil

from app.core.config import get_settings


class PreviewService:
    """Service for generating preview images from PPT files."""
    
    def __init__(self):
        self.settings = get_settings()
        self.data_dir = Path(self.settings.data_dir).resolve()
    
    def _get_preview_dir(self, project_id: str, version: int) -> Path:
        """Get directory for preview images of a version."""
        return self.data_dir / "projects" / project_id / f"v{version}" / "previews"
    
    async def generate_previews(
        self, 
        pptx_path: Path, 
        project_id: str, 
        version: int,
        force: bool = False
    ) -> List[str]:
        """Generate preview images for a PPT file.
        
        Returns list of image filenames (slide_0.png, slide_1.png, etc.)
        """
        preview_dir = self._get_preview_dir(project_id, version)
        
        # Check if previews already exist
        if not force and preview_dir.exists():
            existing = sorted(preview_dir.glob("slide_*.png"))
            if existing:
                return [f.name for f in existing]
        
        preview_dir.mkdir(parents=True, exist_ok=True)
        
        # Try different conversion methods
        success = False
        
        # Method 1: LibreOffice (cross-platform)
        if not success:
            success = await self._convert_with_libreoffice(pptx_path, preview_dir)
        
        # Method 2: Python imaging fallback (placeholder)
        # This would use PIL/Pillow to generate placeholder images
        # In production, you'd want actual PDF/image conversion
        
        if not success:
            # Generate placeholder images as fallback
            success = await self._generate_placeholders(pptx_path, preview_dir)
        
        # Return list of generated images
        images = sorted(preview_dir.glob("slide_*.png"))
        return [f.name for f in images]
    
    async def _convert_with_libreoffice(self, pptx_path: Path, output_dir: Path) -> bool:
        """Convert PPT to images using LibreOffice headless mode."""
        # Check if LibreOffice is available
        libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
        if not libreoffice:
            return False
        
        try:
            # LibreOffice can convert to PDF, then we'd convert PDF to images
            # For simplicity, we'll use a single command approach
            
            # First, convert to PDF
            cmd = [
                libreoffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(output_dir),
                str(pptx_path)
            ]
            
            process = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await process.communicate()
            
            # Find the generated PDF
            pdf_files = list(output_dir.glob("*.pdf"))
            if not pdf_files:
                return False
            
            pdf_path = pdf_files[0]
            
            # Convert PDF to images using pdftoppm or ImageMagick
            pdftoppm = shutil.which("pdftoppm")
            if pdftoppm:
                # pdftoppm from poppler-utils
                cmd = [
                    pdftoppm,
                    "-png",
                    "-r", "150",  # Resolution
                    str(pdf_path),
                    str(output_dir / "slide")
                ]
                process = await asyncio.create_subprocess_exec(
                    *cmd,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE
                )
                await process.communicate()
                
                # Rename files to slide_0.png, slide_1.png, etc.
                for f in sorted(output_dir.glob("slide-*.png")):
                    # slide-1.png -> slide_0.png
                    num = int(f.stem.split("-")[-1]) - 1
                    f.rename(output_dir / f"slide_{num}.png")
            
            # Cleanup PDF
            pdf_path.unlink()
            
            return True
            
        except Exception as e:
            print(f"LibreOffice conversion failed: {e}")
            return False
    
    async def _generate_placeholders(self, pptx_path: Path, output_dir: Path) -> bool:
        """Generate placeholder images when actual conversion isn't available.
        
        This is a fallback that creates simple placeholder images.
        In production, you should install LibreOffice for actual conversion.
        """
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            # Get slide count from pptx
            from pptx import Presentation
            prs = Presentation(str(pptx_path))
            slide_count = len(prs.slides)
            
            for i in range(slide_count):
                # Create placeholder image
                img = Image.new('RGB', (1920, 1080), color='white')
                draw = ImageDraw.Draw(img)
                
                # Draw border
                draw.rectangle([10, 10, 1910, 1070], outline='#cccccc', width=2)
                
                # Draw slide number
                text = f"Slide {i + 1}"
                # Use default font since we might not have custom fonts
                draw.text((960, 540), text, fill='#333333', anchor='mm')
                
                # Draw subtitle
                subtitle = "Preview not available - install LibreOffice for actual previews"
                draw.text((960, 600), subtitle, fill='#999999', anchor='mm')
                
                img.save(output_dir / f"slide_{i}.png", 'PNG')
            
            return True
            
        except ImportError:
            # PIL not available, create empty files as ultimate fallback
            from pptx import Presentation
            prs = Presentation(str(pptx_path))
            slide_count = len(prs.slides)
            
            for i in range(slide_count):
                (output_dir / f"slide_{i}.png").touch()
            
            return True
        except Exception as e:
            print(f"Placeholder generation failed: {e}")
            return False
    
    def get_preview_url(self, project_id: str, version: int, slide_index: int) -> Optional[str]:
        """Get URL path to a preview image."""
        preview_dir = self._get_preview_dir(project_id, version)
        image_path = preview_dir / f"slide_{slide_index}.png"
        
        if image_path.exists():
            # Return relative URL path for the API to serve
            return f"/api/v1/projects/{project_id}/versions/{version}/previews/slide_{slide_index}.png"
        return None
    
    def get_all_preview_urls(self, project_id: str, version: int) -> List[str]:
        """Get URL paths for all preview images of a version."""
        preview_dir = self._get_preview_dir(project_id, version)
        
        if not preview_dir.exists():
            return []
        
        urls = []
        for img in sorted(preview_dir.glob("slide_*.png")):
            slide_index = int(img.stem.split("_")[1])
            urls.append(f"/api/v1/projects/{project_id}/versions/{version}/previews/slide_{slide_index}.png")
        
        return urls
    
    async def delete_previews(self, project_id: str, version: int):
        """Delete preview images for a version."""
        preview_dir = self._get_preview_dir(project_id, version)
        if preview_dir.exists():
            shutil.rmtree(preview_dir)
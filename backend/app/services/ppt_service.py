"""PPT service layer - handles all pptx operations."""

import os
import shutil
import json
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Tuple
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from app.models.ppt import (
    EditInstruction,
    EditInstructionType,
    SlideSummary,
    PPTSummary,
    SlideDetail,
    ProjectMeta,
    VersionMeta,
)
from app.core.config import get_settings


class PPTService:
    """Service for PPT file operations."""
    
    def __init__(self):
        self.settings = get_settings()
        self.data_dir = Path(self.settings.data_dir).resolve()
        self.projects_dir = self.data_dir / "projects"
        self.projects_dir.mkdir(parents=True, exist_ok=True)
    
    def _get_project_dir(self, project_id: str) -> Path:
        """Get project directory path."""
        return self.projects_dir / project_id
    
    def _get_version_path(self, project_id: str, version: int, filename: str = None) -> Path:
        """Get path to a version's file."""
        project_dir = self._get_project_dir(project_id)
        version_dir = project_dir / f"v{version}"
        if filename:
            return version_dir / filename
        return version_dir / "presentation.pptx"
    
    def _get_meta_path(self, project_id: str) -> Path:
        """Get path to project metadata file."""
        return self._get_project_dir(project_id) / "meta.json"
    
    def _load_meta(self, project_id: str) -> dict:
        """Load project metadata."""
        meta_path = self._get_meta_path(project_id)
        if meta_path.exists():
            with open(meta_path, "r", encoding="utf-8") as f:
                return json.load(f)
        return {}
    
    def _save_meta(self, project_id: str, meta: dict):
        """Save project metadata."""
        meta_path = self._get_meta_path(project_id)
        meta_path.parent.mkdir(parents=True, exist_ok=True)
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, indent=2, default=str)
    
    # === Project Management ===
    
    def create_project(self, name: str, topic: str = None) -> Tuple[str, str]:
        """Create a new project. Returns (project_id, initial_pptx_path)."""
        project_id = str(uuid4())[:8]
        project_dir = self._get_project_dir(project_id)
        project_dir.mkdir(parents=True, exist_ok=True)
        
        # Initialize metadata
        now = datetime.now().isoformat()
        meta = {
            "id": project_id,
            "name": name or f"Project {project_id}",
            "created_at": now,
            "updated_at": now,
            "current_version": 0,
            "total_versions": 0,
            "topic": topic,
            "versions": [],
            "pending_edit": None,  # Track unconfirmed edits
        }
        self._save_meta(project_id, meta)
        
        return project_id, str(project_dir)
    
    def get_project_meta(self, project_id: str) -> Optional[ProjectMeta]:
        """Get project metadata."""
        meta = self._load_meta(project_id)
        if not meta:
            return None
        return ProjectMeta(
            id=meta["id"],
            name=meta["name"],
            created_at=datetime.fromisoformat(meta["created_at"]),
            updated_at=datetime.fromisoformat(meta["updated_at"]),
            current_version=meta.get("current_version", 1),
            total_versions=meta.get("total_versions", 1),
            topic=meta.get("topic"),
        )
    
    def list_projects(self) -> List[ProjectMeta]:
        """List all projects."""
        projects = []
        for project_dir in self.projects_dir.iterdir():
            if project_dir.is_dir() and (project_dir / "meta.json").exists():
                meta = self.get_project_meta(project_dir.name)
                if meta:
                    projects.append(meta)
        return sorted(projects, key=lambda p: p.updated_at, reverse=True)
    
    # === Version Management ===
    
    def create_version(self, project_id: str, description: str = None, from_version: int = None) -> int:
        """Create a new version. Returns the new version number."""
        meta = self._load_meta(project_id)
        
        new_version = meta.get("total_versions", 0) + 1
        version_dir = self._get_version_path(project_id, new_version).parent
        version_dir.mkdir(parents=True, exist_ok=True)
        
        # If from_version is specified, copy that version's pptx
        if from_version and from_version > 0:
            src_path = self._get_version_path(project_id, from_version)
            if src_path.exists():
                dst_path = self._get_version_path(project_id, new_version)
                shutil.copy(src_path, dst_path)
        
        # Record version
        version_meta = {
            "version": new_version,
            "created_at": datetime.now().isoformat(),
            "description": description,
            "is_confirmed": True,
        }
        meta["versions"].append(version_meta)
        meta["total_versions"] = new_version
        meta["current_version"] = new_version
        meta["updated_at"] = datetime.now().isoformat()
        self._save_meta(project_id, meta)
        
        return new_version
    
    def get_versions(self, project_id: str) -> List[VersionMeta]:
        """Get all versions of a project."""
        meta = self._load_meta(project_id)
        versions = []
        for v in meta.get("versions", []):
            versions.append(VersionMeta(
                version=v["version"],
                created_at=datetime.fromisoformat(v["created_at"]),
                description=v.get("description"),
                is_confirmed=v.get("is_confirmed", True),
            ))
        return versions
    
    def restore_version(self, project_id: str, target_version: int) -> bool:
        """Restore to a specific version."""
        meta = self._load_meta(project_id)
        
        # Check if target version exists
        version_exists = any(v["version"] == target_version for v in meta.get("versions", []))
        if not version_exists:
            return False
        
        # Create a new version as a copy of target
        new_version = self.create_version(
            project_id, 
            description=f"Restored from version {target_version}",
            from_version=target_version
        )
        
        return True
    
    # === PPT Operations ===
    
    def create_pptx(
        self, 
        project_id: str,
        slides_content: List[dict],
        template_path: str = None
    ) -> Tuple[int, int]:
        """Create a PPT from content. Returns (version, slide_count).
        
        slides_content format:
        [
            {
                "title": "Slide Title",
                "body": ["Point 1", "Point 2"],
                "notes": "Speaker notes (optional)"
            },
            ...
        ]
        """
        # Create first version
        version = self.create_version(project_id, description="Initial creation")
        
        # Create presentation
        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            # Set default slide dimensions (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        
        # Get blank layout
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        
        for i, slide_data in enumerate(slides_content):
            slide = prs.slides.add_slide(blank_layout)
            
            # Add title
            title = slide_data.get("title", "")
            if title:
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = title
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add body content
            body = slide_data.get("body", [])
            if body:
                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(5))
                body_frame = body_box.text_frame
                body_frame.word_wrap = True
                
                for j, point in enumerate(body):
                    if j == 0:
                        p = body_frame.paragraphs[0]
                    else:
                        p = body_frame.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(51, 51, 51)
                    p.space_after = Pt(12)
            
            # Add notes
            notes = slide_data.get("notes", "")
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes
        
        # Save
        pptx_path = self._get_version_path(project_id, version)
        prs.save(str(pptx_path))
        
        return version, len(slides_content)
    
    def execute_edit_instructions(
        self,
        project_id: str,
        version: int,
        instructions: List[EditInstruction]
    ) -> Tuple[int, List[EditInstruction]]:
        """Execute edit instructions and create a pending version.
        Returns (new_version, executed_instructions).
        """
        # Load the source presentation
        pptx_path = self._get_version_path(project_id, version)
        if not pptx_path.exists():
            raise ValueError(f"Version {version} not found for project {project_id}")
        
        prs = Presentation(str(pptx_path))
        
        executed = []
        for instr in instructions:
            try:
                self._execute_single_instruction(prs, instr)
                executed.append(instr)
            except Exception as e:
                # Log error but continue with other instructions
                print(f"Error executing instruction {instr.type}: {e}")
        
        # Create pending version (not yet confirmed)
        meta = self._load_meta(project_id)
        new_version = meta.get("total_versions", 0) + 1
        pending_dir = self._get_version_path(project_id, new_version).parent
        pending_dir.mkdir(parents=True, exist_ok=True)
        
        # Save to pending location
        pending_path = pending_dir / "presentation.pptx"
        prs.save(str(pending_path))
        
        # Track pending edit
        meta["pending_edit"] = {
            "version": new_version,
            "created_at": datetime.now().isoformat(),
            "from_version": version,
        }
        self._save_meta(project_id, meta)
        
        return new_version, executed
    
    def _execute_single_instruction(self, prs: Presentation, instr: EditInstruction):
        """Execute a single edit instruction."""
        slide_index = instr.slide_index
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        
        if instr.type == EditInstructionType.SET_TITLE:
            self._set_title(slide, instr.content)
        elif instr.type == EditInstructionType.SET_BODY:
            self._set_body(slide, instr.content, instr.shape_id)
        elif instr.type == EditInstructionType.SET_BULLETS:
            self._set_bullets(slide, instr.bullets, instr.shape_id)
        elif instr.type == EditInstructionType.ADD_BULLET:
            self._add_bullet(slide, instr.content, instr.shape_id)
        elif instr.type == EditInstructionType.SET_TEXT:
            self._set_text(slide, instr.content, instr.shape_id)
        elif instr.type == EditInstructionType.ADD_TEXT_BOX:
            self._add_text_box(slide, instr.content, instr.position)
        elif instr.type == EditInstructionType.DELETE_SHAPE:
            self._delete_shape(slide, instr.shape_id)
    
    def _find_shape(self, slide, shape_id: str):
        """Find a shape by ID or name."""
        for shape in slide.shapes:
            if shape.name == shape_id:
                return shape
        return None
    
    def _set_title(self, slide, content: str):
        """Set slide title."""
        # Try to find existing title placeholder or text box
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Check if this looks like a title (first text shape or has "title" in name)
                if "title" in shape.name.lower() or shape.top < Inches(1.5):
                    shape.text_frame.paragraphs[0].text = content
                    return
        
        # If no title found, add a new text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_box.text_frame.paragraphs[0].text = content
    
    def _set_body(self, slide, content: str, shape_id: str = None):
        """Set body content."""
        if shape_id:
            shape = self._find_shape(slide, shape_id)
            if shape and shape.has_text_frame:
                shape.text_frame.paragraphs[0].text = content
                return
        
        # Find first body-like shape (below title area)
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top > Inches(1.5):
                shape.text_frame.paragraphs[0].text = content
                return
    
    def _set_bullets(self, slide, bullets: List[str], shape_id: str = None):
        """Set bullet points."""
        if shape_id:
            shape = self._find_shape(slide, shape_id)
            if shape and shape.has_text_frame:
                self._write_bullets(shape.text_frame, bullets)
                return
        
        # Find body shape
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top > Inches(1.5):
                self._write_bullets(shape.text_frame, bullets)
                return
        
        # Add new text box for bullets
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(5))
        self._write_bullets(body_box.text_frame, bullets)
    
    def _write_bullets(self, text_frame, bullets: List[str]):
        """Write bullet points to a text frame."""
        text_frame.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(24)
    
    def _add_bullet(self, slide, content: str, shape_id: str = None):
        """Add a bullet point to existing bullets."""
        if shape_id:
            shape = self._find_shape(slide, shape_id)
            if shape and shape.has_text_frame:
                p = shape.text_frame.add_paragraph()
                p.text = f"• {content}"
                return
        
        # Find body shape and add
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top > Inches(1.5):
                p = shape.text_frame.add_paragraph()
                p.text = f"• {content}"
                return
    
    def _set_text(self, slide, content: str, shape_id: str):
        """Set text of a specific shape."""
        shape = self._find_shape(slide, shape_id)
        if shape and shape.has_text_frame:
            shape.text_frame.paragraphs[0].text = content
    
    def _add_text_box(self, slide, content: str, position: dict = None):
        """Add a new text box."""
        pos = position or {}
        left = Inches(pos.get("left", 0.5))
        top = Inches(pos.get("top", 2))
        width = Inches(pos.get("width", 12.333))
        height = Inches(pos.get("height", 5))
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_box.text_frame.paragraphs[0].text = content
    
    def _delete_shape(self, slide, shape_id: str):
        """Delete a shape by ID."""
        shape = self._find_shape(slide, shape_id)
        if shape:
            sp = shape._element
            sp.getparent().remove(sp)
    
    # === Reading PPT ===
    
    def get_ppt_summary(self, project_id: str, version: int) -> PPTSummary:
        """Get structured summary of a PPT for LLM context."""
        pptx_path = self._get_version_path(project_id, version)
        if not pptx_path.exists():
            raise ValueError(f"Version {version} not found")
        
        prs = Presentation(str(pptx_path))
        slides = []
        
        for i, slide in enumerate(prs.slides):
            summary = self._summarize_slide(slide, i)
            slides.append(summary)
        
        return PPTSummary(
            total_slides=len(prs.slides),
            slides=slides,
            theme=None  # TODO: extract theme info
        )
    
    def _summarize_slide(self, slide, index: int) -> SlideSummary:
        """Summarize a single slide."""
        title = None
        body_text = []
        has_image = False
        has_chart = False
        notes = None
        shape_count = len(slide.shapes)
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                has_chart = True
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                
                # Heuristic: title is the first text shape or has "title" in name
                if title is None and ("title" in shape.name.lower() or shape.top < Inches(1.5)):
                    title = text
                else:
                    body_text.append(text)
        
        # Get notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = notes_text
        
        return SlideSummary(
            index=index,
            title=title,
            body_text="\n".join(body_text) if body_text else None,
            has_image=has_image,
            has_chart=has_chart,
            notes=notes,
            shape_count=shape_count
        )
    
    def get_slide_detail(self, project_id: str, version: int, slide_index: int) -> SlideDetail:
        """Get detailed content of a specific slide."""
        pptx_path = self._get_version_path(project_id, version)
        if not pptx_path.exists():
            raise ValueError(f"Version {version} not found")
        
        prs = Presentation(str(pptx_path))
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {slide_index} out of range")
        
        slide = prs.slides[slide_index]
        shapes = []
        
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            
            if shape.has_text_frame:
                shape_info["text"] = shape.text_frame.text
                shape_info["paragraphs"] = [p.text for p in shape.text_frame.paragraphs]
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info["has_image"] = True
            
            shapes.append(shape_info)
        
        return SlideDetail(index=slide_index, shapes=shapes)
    
    # === Confirm/Cancel Edit ===
    
    def confirm_edit(self, project_id: str, version: int) -> bool:
        """Confirm a pending edit version."""
        meta = self._load_meta(project_id)
        pending = meta.get("pending_edit")
        
        if not pending or pending["version"] != version:
            return False
        
        # Mark as confirmed version
        version_meta = {
            "version": version,
            "created_at": pending["created_at"],
            "description": f"Edit from version {pending['from_version']}",
            "is_confirmed": True,
        }
        meta["versions"].append(version_meta)
        meta["total_versions"] = version
        meta["current_version"] = version
        meta["pending_edit"] = None
        meta["updated_at"] = datetime.now().isoformat()
        self._save_meta(project_id, meta)
        
        return True
    
    def cancel_edit(self, project_id: str, version: int) -> bool:
        """Cancel a pending edit version."""
        meta = self._load_meta(project_id)
        pending = meta.get("pending_edit")
        
        if not pending or pending["version"] != version:
            return False
        
        # Delete pending version files
        pending_dir = self._get_version_path(project_id, version).parent
        if pending_dir.exists():
            shutil.rmtree(pending_dir)
        
        # Clear pending edit
        meta["pending_edit"] = None
        self._save_meta(project_id, meta)
        
        return True
    
    def has_pending_edit(self, project_id: str) -> Optional[int]:
        """Check if there's a pending edit. Returns pending version number or None."""
        meta = self._load_meta(project_id)
        pending = meta.get("pending_edit")
        if pending:
            return pending["version"]
        return None
    
    # === File Paths ===
    
    def get_pptx_path(self, project_id: str, version: int) -> Optional[Path]:
        """Get path to pptx file for a version."""
        path = self._get_version_path(project_id, version)
        return path if path.exists() else None